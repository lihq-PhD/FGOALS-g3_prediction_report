import sys
import cairosvg
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

current_yr=sys.argv[1]
current_mon=sys.argv[2]

prs = Presentation()

pic_box = {
    'Nino3.4': './pic/nino34_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'Nino3': './pic/nino3_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'Nino4': './pic/nino4_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'Nino1+2': './pic/nino12_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'ssta_evolution': './pic/ssta_evolution_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'emi': './pic/emi_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'eio': './pic/eio_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'wio': './pic/wio_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),  
    'dmi': './pic/dmi_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)), 
    'ssta_global': './pic/ssta_global_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)), 
    'hgt850': './pic/hgtuv850_global_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'hgt500': './pic/hgtuv500_global_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'hgt200': './pic/hgtuv200_global_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'precip_land': './pic/precip_land_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'precip_ocean': './pic/precip_ocean_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'temp_land': './pic/temp_land_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'precip_china': './pic/precip_china_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    'temp_china': './pic/temp_china_{}_{}.svg'.format(current_yr, str(current_mon).zfill(2)),
    }

def add_page(page_str):
    page_box = slide.shapes.add_textbox(Inches(9.6), Inches(6.85), Inches(1), Inches(1))      # left, top, width, height
    page_frame = page_box.text_frame
    page_para = page_frame.add_paragraph()
    page_para.text = page_str
    page_para.font.size = Pt(12)
    return slide

def add_title(title_str):
    ct1_box = slide.shapes.add_textbox(Inches(0.1), Inches(0), Inches(8), Inches(0.5))      # left, top, width, height
    ct1_frame = ct1_box.text_frame
    ct1_para = ct1_frame.add_paragraph()
    ct1_para.text = title_str
    ct1_para.font.size = Pt(22)
    ct1_para.font.bold = True
    return slide

def add_pic(pic_str, width, height):
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
        tmp_path = tmp_file.name
    cairosvg.svg2png(url=pic_box[pic_str], write_to=tmp_path, dpi=350)
    pic = slide.shapes.add_picture(tmp_path, (prs.slide_width - width) / 2, Inches(1), width=width, height=height) # , height=Inches(4))
    return slide

# 1. 封面
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1.5))      # left, top, width, height
title_frame = title_box.text_frame
title_para = title_frame.add_paragraph()
title_para.text = "基于FGOALS-g3动力模式的\nENSO预测系统"
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER
title_para.line_spacing = 1.5

subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1.8))      # left, top, width, height
subtitle_frame = subtitle_box.text_frame
subtitle_para = subtitle_frame.add_paragraph()
subtitle_para.text = "{}年{:d}月预报结果".format(current_yr, int(current_mon))
subtitle_para.font.size = Pt(32)
subtitle_para.font.bold = True
subtitle_para.font.color.rgb = RGBColor(255, 0, 0)
subtitle_para.alignment = PP_ALIGN.CENTER

nametitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.8))      # left, top, width, height
nametitle_frame = nametitle_box.text_frame
nametitle_para = nametitle_frame.add_paragraph()
nametitle_para.text = "李昊谦, 陈林, 孙明\n南京信息工程大学"
nametitle_para.font.size = Pt(18)
nametitle_para.alignment = PP_ALIGN.CENTER
nametitle_para.line_spacing = 1.5

timetitle_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1.8))      # left, top, width, height
timetitle_frame = timetitle_box.text_frame
timetitle_para = timetitle_frame.add_paragraph()
timetitle_para.text = datetime.now().strftime("%Y.%m.%d")
timetitle_para.font.size = Pt(18)
timetitle_para.alignment = PP_ALIGN.CENTER


# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("Niño3.4 index")
slide = add_pic('Nino3.4', Inches(6), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("1")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("Niño3 index")
slide = add_pic('Nino3', Inches(6), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("2")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("Niño4 index")
slide = add_pic('Nino4', Inches(6), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("3")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("Niño1+2 index")
slide = add_pic('Nino1+2', Inches(6), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("4")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("SSTA along the equatorial Pacific")
# slide = add_pic('ssta_evolution', Inches(4), None)
with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
    tmp_path = tmp_file.name
cairosvg.svg2png(url=pic_box['ssta_evolution'], write_to=tmp_path, dpi=350)
pic = slide.shapes.add_picture(tmp_path, Inches(1), Inches(1), width=Inches(4), height=None) # , height=Inches(4))

ct2_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("5")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("EMI")
slide = add_pic('emi', Inches(6), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("6")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("DMI")
slide = add_pic('dmi', Inches(6), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("7")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("WIO")
slide = add_pic('wio', Inches(6), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("8")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("EIO")
slide = add_pic('eio', Inches(6), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("9")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("SSTA")
slide = add_pic('ssta_global', Inches(8), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("10")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("HGT&UV 850hPa")
slide = add_pic('hgt850', Inches(8), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("11")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("HGT&UV 500hPa")
slide = add_pic('hgt500', Inches(8), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("12")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("HGT&UV 200hPa")
slide = add_pic('hgt200', Inches(8), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("13")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("Precipitation over Land")
slide = add_pic('precip_land', Inches(8), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("14")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("Precipitation over Ocean")
slide = add_pic('precip_ocean', Inches(8), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("15")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("Temperature over Land")
slide = add_pic('temp_land', Inches(8), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("16")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("Precipitation over China")
slide = add_pic('precip_china', Inches(8), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("17")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("Temperature over China")
slide = add_pic('temp_china', Inches(8), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "暂无"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("18")


prs.save('./realtime_report_{}_{}.pptx'.format(current_yr, str(current_mon).zfill(2)))