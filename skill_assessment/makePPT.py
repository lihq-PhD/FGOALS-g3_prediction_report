import cairosvg
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

prs = Presentation()

pic_box = {
    'ACC&RMSE_Nino3.4': '/mnt/disk3/FG3_hindcast_v1.0/result/FG3_v0.2/f4_ACC_RMSE_nino34.svg',
    'Time_series_Nino3.4': '/mnt/disk3/FG3_hindcast_v1.0/result/FG3_v0.2/f9_Nino34_TCC.svg',
    'PCC_SSTA': '/mnt/disk3/FG3_hindcast_v1.0/result/FG3_v0.2/f8_Nino34_PCC.svg',
    'PCC_thetao': '/mnt/disk3/FG3_hindcast_v1.0/result/FG3_v0.2/f33_thetao_PCC.svg',
    'ACC_SPB': '/mnt/disk3/FG3_hindcast_v1.0/result/FG3_v0.2/f10_Nino34_SPB.svg',
    'FG3_bias': '/mnt/disk3/FG3_hindcast_v1.0/result/FG3_v0.2/f1_model_bias.svg',
    'init_Nino3.4': '/mnt/disk3/FG3_hindcast_v1.0/result/FG3_v0.2/f3_initalization_nino34.svg',
    }

def add_page(page_str):
    page_box = slide.shapes.add_textbox(Inches(9.6), Inches(6.85), Inches(1), Inches(1))      # left, top, width, height
    page_frame = page_box.text_frame
    page_para = page_frame.add_paragraph()
    page_para.text = page_str
    page_para.font.size = Pt(12)
    return slide

def add_title(title_str):
    ct1_box = slide.shapes.add_textbox(Inches(0.1), Inches(0), Inches(8), Inches(0.5))      # left, top, width, height
    ct1_frame = ct1_box.text_frame
    ct1_para = ct1_frame.add_paragraph()
    ct1_para.text = title_str
    ct1_para.font.size = Pt(22)
    ct1_para.font.bold = True
    return slide

def add_pic(pic_str, width, height):
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
        tmp_path = tmp_file.name
    cairosvg.svg2png(url=pic_box[pic_str], write_to=tmp_path, dpi=350)
    pic = slide.shapes.add_picture(tmp_path, (prs.slide_width - width) / 2, Inches(1), width=width, height=height) # , height=Inches(4))
    return slide

# 1. 封面
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(1.5))      # left, top, width, height
title_frame = title_box.text_frame
title_para = title_frame.add_paragraph()
title_para.text = "基于FGOALS-g3动力模式的\nENSO预测系统"
title_para.font.size = Pt(40)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER
title_para.line_spacing = 1.5

subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1.8))      # left, top, width, height
subtitle_frame = subtitle_box.text_frame
subtitle_para = subtitle_frame.add_paragraph()
subtitle_para.text = "预测技巧评估报告"
subtitle_para.font.size = Pt(32)
subtitle_para.font.bold = True
subtitle_para.font.color.rgb = RGBColor(255, 0, 0)
subtitle_para.alignment = PP_ALIGN.CENTER

nametitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.8))      # left, top, width, height
nametitle_frame = nametitle_box.text_frame
nametitle_para = nametitle_frame.add_paragraph()
nametitle_para.text = "李昊谦, 陈林, 孙明\n南京信息工程大学"
nametitle_para.font.size = Pt(18)
nametitle_para.alignment = PP_ALIGN.CENTER
nametitle_para.line_spacing = 1.5

timetitle_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(1.8))      # left, top, width, height
timetitle_frame = timetitle_box.text_frame
timetitle_para = timetitle_frame.add_paragraph()
timetitle_para.text = datetime.now().strftime("%Y.%m.%d")
timetitle_para.font.size = Pt(18)
timetitle_para.alignment = PP_ALIGN.CENTER

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("ACC & RMSE of Niño3.4 index")
slide = add_pic('ACC&RMSE_Nino3.4', Inches(8), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "以1982-2010年为参考期，计算Niño3.4指数的ACC和RMSE：\nFGOALS-g3预测系统能够提前10个月预测ENSO事件，ACC>0.6。"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("1")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("Time series of Niño3.4 index")
slide = add_pic('Time_series_Nino3.4', Inches(5), Inches(5))

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "以1982-2023年为参考期，计算Niño3.4指数的ACC：\nFGOALS-g3预测系统提前8个月预测ENSO事件，ACC>0.6。"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("2")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("PCC of SSTA")
slide = add_pic('PCC_SSTA', Inches(9), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "以1982-2023年为参考期，逐格点计算SSTA的ACC：\nFGOALS-g3预测系统提前8个月在赤道中太平洋具有一定预测技巧，ACC>0.5。"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("3")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("SPB of Niño3.4 index")
slide = add_pic('ACC_SPB', Inches(9), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "以1982-2023年为参考期，计算Niño3.4指数的春季预报障碍（SPB）：\nFGOALS-g3预测系统存在明显的春季预报障碍。"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("4")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("PCC of ocean temperature")
slide = add_pic('PCC_thetao', Inches(8), Inches(5))

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.7), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "以1982-2023年为参考期，逐格点计算5°S-5°N平均的海洋温度ACC：\nFGOALS-g3预测系统在赤道中西太平洋上层300m预测技巧较高，\n东太平洋预测技巧相对较低。"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("5")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("FGOALS-g3 vs OBS")
slide = add_pic('FG3_bias', Inches(6), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "FGOALS-g3预测系统海温气候态偏差。"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("6")

# 2. 内容
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

slide = add_title("Niño3.4 index initialization")
slide = add_pic('init_Nino3.4', Inches(8), None)

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(5.7), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "FGOALS-g3预测系统初始化Niño3.4指数评估。"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

slide = add_page("7")

# 3. 尾页
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

ct1_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.5), Inches(8), Inches(0.5))      # left, top, width, height
ct1_frame = ct1_box.text_frame
ct1_para = ct1_frame.add_paragraph()
ct1_para.text = "观测海表温度资料：HadISST \n https://www.metoffice.gov.uk/hadobs/hadisst/data/download.html"
ct1_para.font.size = Pt(18)
ct1_para.line_spacing = 1.5

ct2_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.5), Inches(8), Inches(0.5))      # left, top, width, height
ct2_frame = ct2_box.text_frame
ct2_para = ct2_frame.add_paragraph()
ct2_para.text = "NMME模式预报结果来源：\n https://www.cpc.ncep.noaa.gov/products/NMME/current/plume.html"
ct2_para.font.size = Pt(18)
ct2_para.line_spacing = 1.5

ct3_box = slide.shapes.add_textbox(Inches(0.3), Inches(3.5), Inches(8), Inches(0.5))      # left, top, width, height
ct3_frame = ct3_box.text_frame
ct3_para = ct3_frame.add_paragraph()
ct3_para.text = "气候态参考期：1982-2010年 或 1991-2020年"
ct3_para.font.size = Pt(18)
ct3_para.line_spacing = 1.5

ct4_box = slide.shapes.add_textbox(Inches(0.3), Inches(4), Inches(8), Inches(0.5))      # left, top, width, height
ct4_frame = ct4_box.text_frame
ct4_para = ct4_frame.add_paragraph()
ct4_para.text = "初始化：SST-Nudging（全场SST初始化、异常场SSTA初始化）"
ct4_para.font.size = Pt(18)
ct4_para.line_spacing = 1.5

ct5_box = slide.shapes.add_textbox(Inches(0.3), Inches(4.5), Inches(8), Inches(0.5))      # left, top, width, height
ct5_frame = ct5_box.text_frame
ct5_para = ct5_frame.add_paragraph()
ct5_para.text = "成员数：8个成员"
ct5_para.font.size = Pt(18)
ct5_para.line_spacing = 1.5

ct6_box = slide.shapes.add_textbox(Inches(0.3), Inches(6), Inches(8), Inches(0.5))      # left, top, width, height
ct6_frame = ct6_box.text_frame
ct6_para = ct6_frame.add_paragraph()
ct6_para.text = "如有疑问，请联系：李昊谦 lihq_enso@163.com"
ct6_para.font.size = Pt(18)
ct6_para.line_spacing = 1.5

slide = add_page("8")

prs.save('./prediction_skill_report.pptx')